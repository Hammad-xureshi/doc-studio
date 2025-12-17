import tempfile
import streamlit as st
import google.generativeai as genai
from pathlib import Path
import os
import hashlib
from typing import Dict, List
from datetime import datetime
import time
import json

# Document Processing
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

# Vector DB
import chromadb
from chromadb.config import Settings as ChromaSettings

# Data & Visualization
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go

from dotenv import load_dotenv

load_dotenv()

# ============================================
# ⚙️ CONFIGURATION
# ============================================

API_KEY = os.getenv("GEMINI_API_KEY", "")
CREATOR = "Hammad Naeem"
VERSION = "3.0 LUXE"

# ============================================
# 🎨 PAGE CONFIG
# ============================================

st.set_page_config(
    page_title=f"✨ Doc Studio - By {CREATOR}",
    page_icon="✨",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ============================================
# 🌸 PINTEREST-STYLE FUTURISTIC THEME
# ============================================

st.markdown("""
<style>
    /* ===== GOOGLE FONTS ===== */
    @import url('https://fonts.googleapis.com/css2?family=Poppins:wght@300;400;500;600;700;800&family=Space+Grotesk:wght@300;400;500;600;700&display=swap');
    
    /* ===== ROOT VARIABLES ===== */
    :root {
        --bg-primary: #0f0f1a;
        --bg-secondary: #1a1a2e;
        --bg-card: rgba(255, 255, 255, 0.03);
        --glass-bg: rgba(255, 255, 255, 0.05);
        --glass-border: rgba(255, 255, 255, 0.1);
        
        --neon-pink: #ff6b9d;
        --neon-purple: #c56cf0;
        --neon-blue: #70a1ff;
        --neon-cyan: #7bed9f;
        --neon-orange: #ffa502;
        --neon-mint: #a3f7bf;
        
        --gradient-1: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        --gradient-2: linear-gradient(135deg, #f093fb 0%, #f5576c 100%);
        --gradient-3: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%);
        --gradient-4: linear-gradient(135deg, #43e97b 0%, #38f9d7 100%);
        --gradient-5: linear-gradient(135deg, #fa709a 0%, #fee140 100%);
        --gradient-hero: linear-gradient(135deg, #667eea 0%, #764ba2 50%, #f093fb 100%);
        
        --text-primary: #ffffff;
        --text-secondary: rgba(255, 255, 255, 0.7);
        --text-muted: rgba(255, 255, 255, 0.5);
        
        --shadow-soft: 0 8px 32px rgba(0, 0, 0, 0.3);
        --shadow-glow: 0 0 40px rgba(102, 126, 234, 0.3);
        --shadow-neon: 0 0 20px rgba(255, 107, 157, 0.4);
        
        --radius-sm: 12px;
        --radius-md: 20px;
        --radius-lg: 30px;
        --radius-xl: 40px;
    }
    
    /* ===== HIDE STREAMLIT DEFAULTS ===== */
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}
    .stDeployButton {display: none;}
    
    /* ===== GLOBAL STYLES ===== */
    * {
        font-family: 'Poppins', sans-serif;
        transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1);
    }
    
    /* ===== MAIN BACKGROUND ===== */
    .stApp {
        background: var(--bg-primary);
        background-image: 
            radial-gradient(ellipse at 20% 20%, rgba(102, 126, 234, 0.15) 0%, transparent 50%),
            radial-gradient(ellipse at 80% 80%, rgba(240, 147, 251, 0.15) 0%, transparent 50%),
            radial-gradient(ellipse at 40% 60%, rgba(79, 172, 254, 0.1) 0%, transparent 40%);
        min-height: 100vh;
    }
    
    /* ===== FLOATING PARTICLES ANIMATION ===== */
    .stApp::before {
        content: "";
        position: fixed;
        top: 0;
        left: 0;
        width: 100%;
        height: 100%;
        background-image: 
            radial-gradient(2px 2px at 20px 30px, rgba(255,255,255,0.1), transparent),
            radial-gradient(2px 2px at 40px 70px, rgba(255,255,255,0.08), transparent),
            radial-gradient(1px 1px at 90px 40px, rgba(255,255,255,0.1), transparent),
            radial-gradient(2px 2px at 130px 80px, rgba(255,255,255,0.06), transparent);
        background-size: 200px 200px;
        animation: floatParticles 20s linear infinite;
        pointer-events: none;
        z-index: 0;
    }
    
    @keyframes floatParticles {
        0% { transform: translateY(0); }
        100% { transform: translateY(-200px); }
    }
    
    /* ===== SMOOTH FADE IN ANIMATION ===== */
    @keyframes fadeInUp {
        from {
            opacity: 0;
            transform: translateY(30px);
        }
        to {
            opacity: 1;
            transform: translateY(0);
        }
    }
    
    @keyframes fadeIn {
        from { opacity: 0; }
        to { opacity: 1; }
    }
    
    @keyframes slideIn {
        from {
            opacity: 0;
            transform: translateX(-20px);
        }
        to {
            opacity: 1;
            transform: translateX(0);
        }
    }
    
    @keyframes pulse {
        0%, 100% { transform: scale(1); }
        50% { transform: scale(1.02); }
    }
    
    @keyframes shimmer {
        0% { background-position: -200% 0; }
        100% { background-position: 200% 0; }
    }
    
    @keyframes glow {
        0%, 100% { box-shadow: 0 0 20px rgba(102, 126, 234, 0.4); }
        50% { box-shadow: 0 0 40px rgba(240, 147, 251, 0.6); }
    }
    
    /* ===== SIDEBAR STYLING ===== */
    [data-testid="stSidebar"] {
        background: linear-gradient(180deg, rgba(26, 26, 46, 0.95) 0%, rgba(15, 15, 26, 0.98) 100%);
        backdrop-filter: blur(20px);
        border-right: 1px solid var(--glass-border);
        box-shadow: 4px 0 30px rgba(0, 0, 0, 0.3);
    }
    
    [data-testid="stSidebar"] > div:first-child {
        padding-top: 0;
    }
    
    [data-testid="stSidebar"] .stMarkdown {
        color: var(--text-primary);
    }
    
    [data-testid="stSidebar"] label {
        color: var(--text-secondary) !important;
        font-weight: 500;
    }
    
    /* ===== GLASSMORPHISM CARDS ===== */
    .glass-card {
        background: var(--glass-bg);
        backdrop-filter: blur(20px);
        -webkit-backdrop-filter: blur(20px);
        border-radius: var(--radius-lg);
        border: 1px solid var(--glass-border);
        padding: 30px;
        margin: 20px 0;
        box-shadow: var(--shadow-soft);
        animation: fadeInUp 0.6s ease-out;
        position: relative;
        overflow: hidden;
    }
    
    .glass-card::before {
        content: "";
        position: absolute;
        top: 0;
        left: -100%;
        width: 100%;
        height: 100%;
        background: linear-gradient(90deg, transparent, rgba(255,255,255,0.05), transparent);
        transition: left 0.5s ease;
    }
    
    .glass-card:hover::before {
        left: 100%;
    }
    
    .glass-card:hover {
        transform: translateY(-5px);
        box-shadow: var(--shadow-soft), var(--shadow-glow);
        border-color: rgba(102, 126, 234, 0.3);
    }
    
    /* ===== HERO HEADER ===== */
    .hero-header {
        font-family: 'Space Grotesk', sans-serif;
        font-size: 3.5rem;
        font-weight: 700;
        text-align: center;
        background: var(--gradient-hero);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        background-clip: text;
        margin-bottom: 10px;
        animation: fadeInUp 0.8s ease-out;
        letter-spacing: -1px;
    }
    
    .hero-sub {
        text-align: center;
        color: var(--text-secondary);
        font-size: 1.2rem;
        font-weight: 400;
        margin-bottom: 10px;
        animation: fadeInUp 0.8s ease-out 0.1s both;
    }
    
    .hero-creator {
        text-align: center;
        font-size: 1rem;
        font-weight: 600;
        margin-bottom: 30px;
        animation: fadeInUp 0.8s ease-out 0.2s both;
        background: var(--gradient-2);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
    }
    
    /* ===== CHAT MESSAGES ===== */
    .user-message {
        background: var(--gradient-1);
        color: white;
        padding: 20px 25px;
        border-radius: var(--radius-lg) var(--radius-lg) 8px var(--radius-lg);
        margin: 16px 0;
        margin-left: 15%;
        box-shadow: var(--shadow-soft);
        animation: slideIn 0.4s ease-out;
        font-weight: 500;
        line-height: 1.6;
    }
    
    .ai-message {
        background: var(--glass-bg);
        backdrop-filter: blur(20px);
        color: var(--text-primary);
        padding: 25px 30px;
        border-radius: var(--radius-lg) var(--radius-lg) var(--radius-lg) 8px;
        margin: 16px 0;
        margin-right: 15%;
        border: 1px solid var(--glass-border);
        box-shadow: var(--shadow-soft);
        animation: slideIn 0.4s ease-out;
        line-height: 1.8;
    }
    
    .ai-message strong {
        color: var(--neon-pink);
    }
    
    /* ===== SOURCE CITATION ===== */
    .source-tag {
        display: inline-flex;
        align-items: center;
        gap: 8px;
        background: linear-gradient(135deg, rgba(102, 126, 234, 0.2), rgba(118, 75, 162, 0.2));
        border: 1px solid rgba(102, 126, 234, 0.3);
        padding: 10px 18px;
        border-radius: 50px;
        margin: 8px 8px 8px 0;
        font-size: 0.85rem;
        color: var(--neon-blue);
        font-weight: 500;
        transition: all 0.3s ease;
    }
    
    .source-tag:hover {
        background: linear-gradient(135deg, rgba(102, 126, 234, 0.3), rgba(118, 75, 162, 0.3));
        transform: translateY(-2px);
        box-shadow: 0 5px 20px rgba(102, 126, 234, 0.3);
    }
    
    /* ===== BUTTONS ===== */
    .stButton > button {
        background: var(--gradient-1);
        color: white;
        border: none;
        border-radius: 50px;
        padding: 14px 35px;
        font-weight: 600;
        font-size: 1rem;
        letter-spacing: 0.5px;
        box-shadow: 0 10px 30px rgba(102, 126, 234, 0.4);
        transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1);
        position: relative;
        overflow: hidden;
    }
    
    .stButton > button::before {
        content: "";
        position: absolute;
        top: 0;
        left: -100%;
        width: 100%;
        height: 100%;
        background: linear-gradient(90deg, transparent, rgba(255,255,255,0.2), transparent);
        transition: left 0.5s ease;
    }
    
    .stButton > button:hover::before {
        left: 100%;
    }
    
    .stButton > button:hover {
        transform: translateY(-3px) scale(1.02);
        box-shadow: 0 15px 40px rgba(102, 126, 234, 0.5);
    }
    
    .stButton > button:active {
        transform: translateY(-1px) scale(0.98);
    }
    
    /* ===== INPUT FIELDS ===== */
    .stTextInput > div > div > input,
    .stTextArea > div > div > textarea {
        background: var(--glass-bg);
        backdrop-filter: blur(10px);
        border: 2px solid var(--glass-border);
        border-radius: var(--radius-md);
        color: var(--text-primary);
        padding: 16px 22px;
        font-size: 1rem;
        font-weight: 500;
        transition: all 0.3s ease;
    }
    
    .stTextInput > div > div > input:focus,
    .stTextArea > div > div > textarea:focus {
        border-color: var(--neon-purple);
        box-shadow: 0 0 0 4px rgba(197, 108, 240, 0.15), var(--shadow-neon);
        background: rgba(197, 108, 240, 0.05);
    }
    
    .stTextInput > div > div > input::placeholder,
    .stTextArea > div > div > textarea::placeholder {
        color: var(--text-muted);
    }
    
    /* ===== FILE UPLOADER ===== */
    [data-testid="stFileUploader"] {
        background: var(--glass-bg);
        border: 2px dashed var(--glass-border);
        border-radius: var(--radius-lg);
        padding: 30px;
        transition: all 0.3s ease;
    }
    
    [data-testid="stFileUploader"]:hover {
        border-color: var(--neon-purple);
        background: rgba(197, 108, 240, 0.05);
        box-shadow: 0 0 30px rgba(197, 108, 240, 0.2);
    }
    
    [data-testid="stFileUploader"] label {
        color: var(--text-secondary) !important;
    }
    
    /* ===== TABS ===== */
    .stTabs [data-baseweb="tab-list"] {
        gap: 10px;
        background: transparent;
        padding: 10px;
        border-radius: var(--radius-lg);
    }
    
    .stTabs [data-baseweb="tab"] {
        background: var(--glass-bg);
        backdrop-filter: blur(10px);
        border-radius: 50px;
        padding: 14px 28px;
        color: var(--text-secondary);
        border: 1px solid var(--glass-border);
        font-weight: 600;
        font-size: 0.95rem;
        transition: all 0.3s ease;
    }
    
    .stTabs [data-baseweb="tab"]:hover {
        background: rgba(102, 126, 234, 0.1);
        color: var(--text-primary);
        border-color: rgba(102, 126, 234, 0.3);
    }
    
    .stTabs [aria-selected="true"] {
        background: var(--gradient-1) !important;
        color: white !important;
        border: none !important;
        box-shadow: 0 10px 30px rgba(102, 126, 234, 0.4);
    }
    
    /* ===== METRICS CARDS ===== */
    .metric-box {
        background: var(--glass-bg);
        backdrop-filter: blur(20px);
        border-radius: var(--radius-lg);
        padding: 25px;
        text-align: center;
        border: 1px solid var(--glass-border);
        transition: all 0.3s ease;
        animation: fadeInUp 0.6s ease-out;
    }
    
    .metric-box:hover {
        transform: translateY(-8px) scale(1.02);
        box-shadow: var(--shadow-glow);
        border-color: rgba(102, 126, 234, 0.4);
    }
    
    .metric-number {
        font-family: 'Space Grotesk', sans-serif;
        font-size: 2.8rem;
        font-weight: 700;
        background: var(--gradient-1);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        line-height: 1.2;
    }
    
    .metric-title {
        color: var(--text-secondary);
        font-size: 0.9rem;
        font-weight: 600;
        margin-top: 8px;
        text-transform: uppercase;
        letter-spacing: 1px;
    }
    
    /* ===== FEATURE CARDS ===== */
    .feature-box {
        background: var(--glass-bg);
        backdrop-filter: blur(20px);
        border-radius: var(--radius-lg);
        padding: 35px 25px;
        text-align: center;
        border: 1px solid var(--glass-border);
        transition: all 0.4s cubic-bezier(0.4, 0, 0.2, 1);
        animation: fadeInUp 0.6s ease-out;
        position: relative;
        overflow: hidden;
    }
    
    .feature-box::after {
        content: "";
        position: absolute;
        bottom: 0;
        left: 0;
        width: 100%;
        height: 3px;
        background: var(--gradient-1);
        transform: scaleX(0);
        transition: transform 0.4s ease;
    }
    
    .feature-box:hover::after {
        transform: scaleX(1);
    }
    
    .feature-box:hover {
        transform: translateY(-10px);
        box-shadow: var(--shadow-soft), var(--shadow-glow);
        border-color: rgba(102, 126, 234, 0.3);
    }
    
    .feature-icon {
        font-size: 3.5rem;
        margin-bottom: 20px;
        display: block;
    }
    
    .feature-title {
        color: var(--text-primary);
        font-size: 1.2rem;
        font-weight: 700;
        margin-bottom: 12px;
    }
    
    .feature-desc {
        color: var(--text-secondary);
        font-size: 0.9rem;
        line-height: 1.6;
    }
    
    /* ===== SELECTBOX ===== */
    .stSelectbox > div > div {
        background: var(--glass-bg);
        border: 2px solid var(--glass-border);
        border-radius: var(--radius-md);
        color: var(--text-primary);
    }
    
    .stSelectbox > div > div:hover {
        border-color: var(--neon-purple);
    }
    
    /* ===== SLIDER ===== */
    .stSlider > div > div > div {
        background: var(--gradient-1);
    }
    
    /* ===== EXPANDER ===== */
    .streamlit-expanderHeader {
        background: var(--glass-bg);
        border-radius: var(--radius-md);
        border: 1px solid var(--glass-border);
        color: var(--text-primary) !important;
        font-weight: 600;
    }
    
    .streamlit-expanderHeader:hover {
        border-color: var(--neon-purple);
        background: rgba(197, 108, 240, 0.05);
    }
    
    /* ===== SUCCESS/ERROR/WARNING ===== */
    .success-box {
        background: linear-gradient(135deg, rgba(67, 233, 123, 0.15), rgba(56, 249, 215, 0.1));
        border-left: 4px solid var(--neon-cyan);
        padding: 18px 24px;
        border-radius: 0 var(--radius-md) var(--radius-md) 0;
        color: var(--neon-cyan);
        font-weight: 600;
        animation: slideIn 0.4s ease-out;
    }
    
    .error-box {
        background: linear-gradient(135deg, rgba(255, 107, 157, 0.15), rgba(245, 87, 108, 0.1));
        border-left: 4px solid var(--neon-pink);
        padding: 18px 24px;
        border-radius: 0 var(--radius-md) var(--radius-md) 0;
        color: var(--neon-pink);
        font-weight: 600;
        animation: slideIn 0.4s ease-out;
    }
    
    .warning-box {
        background: linear-gradient(135deg, rgba(255, 165, 2, 0.15), rgba(254, 225, 64, 0.1));
        border-left: 4px solid var(--neon-orange);
        padding: 18px 24px;
        border-radius: 0 var(--radius-md) var(--radius-md) 0;
        color: var(--neon-orange);
        font-weight: 600;
        animation: slideIn 0.4s ease-out;
    }
    
    .info-box {
        background: linear-gradient(135deg, rgba(112, 161, 255, 0.15), rgba(79, 172, 254, 0.1));
        border-left: 4px solid var(--neon-blue);
        padding: 18px 24px;
        border-radius: 0 var(--radius-md) var(--radius-md) 0;
        color: var(--neon-blue);
        font-weight: 500;
        animation: slideIn 0.4s ease-out;
        line-height: 1.6;
    }
    
    /* ===== DOWNLOAD BUTTON ===== */
    .stDownloadButton > button {
        background: var(--gradient-4) !important;
        box-shadow: 0 10px 30px rgba(67, 233, 123, 0.4);
    }
    
    .stDownloadButton > button:hover {
        box-shadow: 0 15px 40px rgba(67, 233, 123, 0.5);
    }
    
    /* ===== SCROLLBAR ===== */
    ::-webkit-scrollbar {
        width: 10px;
        height: 10px;
    }
    
    ::-webkit-scrollbar-track {
        background: var(--bg-secondary);
        border-radius: 10px;
    }
    
    ::-webkit-scrollbar-thumb {
        background: var(--gradient-1);
        border-radius: 10px;
    }
    
    ::-webkit-scrollbar-thumb:hover {
        background: var(--gradient-2);
    }
    
    /* ===== DIVIDER ===== */
    .fancy-divider {
        height: 2px;
        background: linear-gradient(90deg, transparent, var(--glass-border), transparent);
        margin: 30px 0;
        border: none;
    }
    
    /* ===== LOGIN CARD ===== */
    .login-container {
        max-width: 450px;
        margin: 60px auto;
        padding: 50px 40px;
        background: var(--glass-bg);
        backdrop-filter: blur(30px);
        border-radius: var(--radius-xl);
        border: 1px solid var(--glass-border);
        box-shadow: var(--shadow-soft), var(--shadow-glow);
        animation: fadeInUp 0.8s ease-out;
    }
    
    .login-title {
        font-family: 'Space Grotesk', sans-serif;
        font-size: 2rem;
        font-weight: 700;
        text-align: center;
        background: var(--gradient-1);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        margin-bottom: 10px;
    }
    
    .login-subtitle {
        text-align: center;
        color: var(--text-secondary);
        font-size: 1rem;
        margin-bottom: 30px;
    }
    
    /* ===== BADGE ===== */
    .badge {
        display: inline-flex;
        align-items: center;
        gap: 6px;
        padding: 6px 14px;
        border-radius: 50px;
        font-size: 0.8rem;
        font-weight: 600;
    }
    
    .badge-success {
        background: linear-gradient(135deg, rgba(67, 233, 123, 0.2), rgba(56, 249, 215, 0.15));
        color: var(--neon-cyan);
        border: 1px solid rgba(67, 233, 123, 0.3);
    }
    
    .badge-error {
        background: linear-gradient(135deg, rgba(255, 107, 157, 0.2), rgba(245, 87, 108, 0.15));
        color: var(--neon-pink);
        border: 1px solid rgba(255, 107, 157, 0.3);
    }
    
    /* ===== FOOTER ===== */
    .footer-container {
        background: var(--glass-bg);
        backdrop-filter: blur(20px);
        border-radius: var(--radius-lg);
        padding: 40px;
        text-align: center;
        border: 1px solid var(--glass-border);
        margin-top: 50px;
        animation: fadeInUp 0.8s ease-out;
    }
    
    .footer-title {
        font-family: 'Space Grotesk', sans-serif;
        font-size: 1.8rem;
        font-weight: 700;
        background: var(--gradient-1);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        margin-bottom: 15px;
    }
    
    .footer-creator {
        font-size: 1.1rem;
        font-weight: 600;
        background: var(--gradient-2);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        margin: 15px 0;
    }
    
    .footer-text {
        color: var(--text-secondary);
        font-size: 0.9rem;
        line-height: 1.8;
    }
    
    /* ===== RADIO BUTTONS ===== */
    .stRadio > div {
        background: var(--glass-bg);
        padding: 15px 20px;
        border-radius: var(--radius-md);
        border: 1px solid var(--glass-border);
    }
    
    .stRadio label {
        color: var(--text-primary) !important;
    }
    
    /* ===== CONTENT BOX ===== */
    .content-box {
        background: var(--glass-bg);
        backdrop-filter: blur(20px);
        border-radius: var(--radius-lg);
        padding: 30px;
        border: 1px solid var(--glass-border);
        margin: 20px 0;
        animation: fadeInUp 0.6s ease-out;
        line-height: 1.8;
        color: var(--text-primary);
    }
    
    .content-box h1, .content-box h2, .content-box h3 {
        color: var(--neon-blue);
        margin-bottom: 15px;
    }
    
    .content-box ul, .content-box ol {
        margin-left: 20px;
        color: var(--text-secondary);
    }
    
    .content-box strong {
        color: var(--neon-pink);
    }
    
    /* ===== GLOW EFFECT ===== */
    .glow-text {
        text-shadow: 0 0 10px currentColor;
    }
    
    /* ===== MOBILE RESPONSIVE ===== */
    @media (max-width: 768px) {
        .hero-header {
            font-size: 2.2rem;
        }
        
        .hero-sub {
            font-size: 1rem;
        }
        
        .user-message, .ai-message {
            margin-left: 5%;
            margin-right: 5%;
            padding: 16px 20px;
        }
        
        .glass-card {
            padding: 20px;
        }
        
        .metric-number {
            font-size: 2rem;
        }
        
        .feature-box {
            padding: 25px 20px;
        }
        
        .login-container {
            margin: 30px 15px;
            padding: 35px 25px;
        }
    }
</style>
""", unsafe_allow_html=True)


# ============================================
# 🔧 HELPER FUNCTIONS
# ============================================

def show_success(message):
    st.markdown(f'<div class="success-box">✅ {message}</div>', unsafe_allow_html=True)

def show_error(message):
    st.markdown(f'<div class="error-box">❌ {message}</div>', unsafe_allow_html=True)

def show_warning(message):
    st.markdown(f'<div class="warning-box">⚠️ {message}</div>', unsafe_allow_html=True)

def show_info(message):
    st.markdown(f'<div class="info-box">💡 {message}</div>', unsafe_allow_html=True)

def add_watermark(content: str) -> str:
    """Add creator watermark"""
    header = f"""
{'='*60}
✨ AI DOCUMENT STUDIO {VERSION}
📌 Created by: {CREATOR}
📅 Generated: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}
{'='*60}

"""
    footer = f"""

{'='*60}
© {datetime.now().year} {CREATOR} - All Rights Reserved
{'='*60}
"""
    return header + content + footer


def find_model(api_key: str):
    """
    Try to connect to Gemini and return a working model.
    Returns: (model_name or None, error_message or None)
    """
    if not api_key or len(api_key.strip()) < 10:
        return None, "Invalid API key format"
    
    api_key = api_key.strip()

    try:
        genai.configure(api_key=api_key)
    except Exception as e:
        return None, f"Failed to configure API: {str(e)[:120]}"

    # Sab se pehle 1 stable model test karo
    preferred_models = [
        "gemini-1.5-flash",          # most common & recommended
        "gemini-1.5-flash-latest",
        "gemini-1.5-pro",
        "gemini-1.5-pro-latest",
    ]

    last_error = ""

    for m in preferred_models:
        try:
            model = genai.GenerativeModel(m)
            resp = model.generate_content("test", generation_config={"max_output_tokens": 5})
            if resp and getattr(resp, "text", "").strip():
                return m, None
        except Exception as e:
            last_error = str(e)
            continue

    # Agar upar se kuch na mila, to API se models list karke dekh lo
    try:
        available = []
        for mdl in genai.list_models():
            if "generateContent" in str(mdl.supported_generation_methods):
                name = mdl.name.replace("models/", "")
                available.append(name)
        
        for name in available:
            try:
                model = genai.GenerativeModel(name)
                resp = model.generate_content("test", generation_config={"max_output_tokens": 5})
                if resp and getattr(resp, "text", "").strip():
                    return name, None
            except Exception as e:
                last_error = str(e)
                continue

        if available:
            return None, f"No working model. Available: {', '.join(available[:3])}"
        else:
            return None, "No text generation models available for this API key."

    except Exception as e:
        if not last_error:
            last_error = str(e)
        return None, f"Model discovery failed: {last_error[:120]}"


# ============================================
# 📄 DOCUMENT PROCESSOR
# ============================================

class DocumentProcessor:
    def process(self, path: str) -> Dict:
        """Process uploaded document"""
        p = Path(path)
        ext = p.suffix.lower()
        
        try:
            if ext == ".pdf":
                content, pages = self._read_pdf(path)
            elif ext == ".docx":
                content, pages = self._read_docx(path)
            elif ext == ".xlsx":
                content, pages = self._read_excel(path)
            elif ext == ".pptx":
                content, pages = self._read_pptx(path)
            elif ext == ".txt":
                content, pages = self._read_txt(path)
            else:
                raise ValueError(f"Unsupported: {ext}")
            
            doc_id = hashlib.md5(f"{p.name}{datetime.now()}".encode()).hexdigest()[:16]
            
            return {
                "id": doc_id,
                "name": p.name,
                "type": ext[1:].upper(),
                "content": content,
                "chunks": self._create_chunks(pages, p.name, doc_id),
                "pages": len(pages),
                "size": p.stat().st_size,
                "uploaded": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                "word_count": len(content.split())
            }
        except Exception as e:
            raise Exception(f"Error: {str(e)}")
    
    def _read_pdf(self, path):
        reader = PdfReader(path)
        pages = {}
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            pages[i] = text.strip()
        content = "\n\n".join([f"[Page {k}]\n{v}" for k, v in pages.items() if v])
        return content, pages
    
    def _read_docx(self, path):
        doc = DocxDocument(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n\n".join(paragraphs)
        return text, {1: text}
    
    def _read_excel(self, path):
        wb = load_workbook(path, data_only=True)
        pages = {}
        for i, name in enumerate(wb.sheetnames, 1):
            sheet = wb[name]
            rows = []
            for row in sheet.iter_rows(values_only=True):
                r = " | ".join(str(c) if c else "" for c in row)
                if r.strip(" |"):
                    rows.append(r)
            pages[i] = f"[Sheet: {name}]\n" + "\n".join(rows)
        return "\n\n".join(pages.values()), pages
    
    def _read_pptx(self, path):
        prs = Presentation(path)
        pages = {}
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            pages[i] = "\n".join(texts)
        content = "\n\n".join([f"[Slide {k}]\n{v}" for k, v in pages.items()])
        return content, pages
    
    def _read_txt(self, path):
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        return text, {1: text}
    
    def _create_chunks(self, pages, filename, doc_id):
        chunks = []
        chunk_size = 500
        overlap = 100
        
        for page_num, text in pages.items():
            if not text.strip():
                continue
            words = text.split()
            
            for i in range(0, len(words), chunk_size - overlap):
                chunk_words = words[i:i + chunk_size]
                chunk_text = " ".join(chunk_words)
                
                if len(chunk_text) > 80:
                    chunks.append({
                        "id": f"{doc_id}_p{page_num}_c{i}",
                        "text": chunk_text,
                        "page": page_num,
                        "file": filename
                    })
        return chunks


# ============================================
# 🗄️ VECTOR DATABASE
# ============================================

class VectorStore:
    def __init__(self):
        # Memory-based DB setup (Vercel ke liye best)
        os.environ['ANONYMIZED_TELEMETRY'] = 'False'
        
        self.client = chromadb.EphemeralClient(
            settings=ChromaSettings(
                anonymized_telemetry=False,
                allow_reset=True
            )
        )
        self.collection = self.client.get_or_create_collection(
            name="docs",
            metadata={"hnsw:space": "cosine"}
        )
    
    
    def _embed(self, text, task="retrieval_document"):
        try:
            result = genai.embed_content(
                model="models/embedding-001",
                content=text[:8000],
                task_type=task
            )
            return result['embedding']
        except:
            return [0.0] * 768
    
    def add_document(self, doc_id, chunks):
        if not chunks:
            return
        
        ids, embeddings, documents, metadatas = [], [], [], []
        
        for chunk in chunks:
            try:
                ids.append(f"{doc_id}_{chunk['id']}")
                embeddings.append(self._embed(chunk['text']))
                documents.append(chunk['text'])
                metadatas.append({
                    "doc_id": doc_id,
                    "page": str(chunk['page']),
                    "file": chunk['file']
                })
            except:
                continue
        
        if ids:
            try:
                self.collection.add(
                    ids=ids,
                    embeddings=embeddings,
                    documents=documents,
                    metadatas=metadatas
                )
            except Exception as e:
                pass
    
    def search(self, query, k=5):
        try:
            count = self.collection.count()
            if count == 0:
                return []
            
            embedding = self._embed(query, "retrieval_query")
            results = self.collection.query(
                query_embeddings=[embedding],
                n_results=min(k, count)
            )
            
            output = []
            if results['documents'] and results['documents'][0]:
                for i, doc in enumerate(results['documents'][0]):
                    output.append({
                        "text": doc,
                        "meta": results['metadatas'][0][i] if results['metadatas'] else {},
                        "distance": results['distances'][0][i] if results.get('distances') else 0
                    })
            return sorted(output, key=lambda x: x.get('distance', 999))
        except:
            return []
    
    def delete_document(self, doc_id):
        try:
            results = self.collection.get(where={"doc_id": doc_id})
            if results['ids']:
                self.collection.delete(ids=results['ids'])
        except:
            pass
    
    def count(self):
        try:
            return self.collection.count()
        except:
            return 0


# ============================================
# 🤖 AI ASSISTANT
# ============================================

class AIAssistant:
    def __init__(self, model_name: str):
        self.model_name = model_name
        self.model = genai.GenerativeModel(
            model_name,
            generation_config={
                "temperature": 0.7,
                "top_p": 0.9,
                "max_output_tokens": 2048,
            },
            safety_settings={
                "harassment": "block_none",
                "hate": "block_none",
                "sex": "block_none",
                "danger": "block_none",
            }
        )

    def _extract_text(self, response):
        """Safely extract text from Gemini response"""
        if not response:
            return ""
        # Latest google-generativeai responses normally have .text
        if hasattr(response, "text") and response.text:
            return response.text
        # fallback to candidates / parts (rarely needed)
        try:
            parts = []
            for cand in getattr(response, "candidates", []):
                for part in getattr(cand, "content", {}).get("parts", []):
                    if isinstance(part, dict) and "text" in part:
                        parts.append(part["text"])
            return "\n".join(parts)
        except:
            return ""

    def generate(self, prompt: str) -> str:
        """Core generate function used everywhere"""
        for attempt in range(3):
            try:
                response = self.model.generate_content(prompt)
                text = self._extract_text(response)
                if text.strip():
                    return text
                else:
                    return "Model did not return any text."
            except Exception as e:
                err = str(e).lower()
                if "rate" in err or "quota" in err or "429" in err:
                    time.sleep(1 + attempt)
                    continue
                if "safety" in err:
                    return "Response blocked by safety filters from Gemini."
                # last attempt: return error
                if attempt == 2:
                    return f"Error from model: {str(e)[:200]}"
                time.sleep(1)
        return "Failed after multiple attempts."
    
    def answer_question(self, question, context, language="en"):
        """Answer question from context"""
        context_text = ""
        for i, c in enumerate(context, 1):
            meta = c.get("meta", {})
            context_text += f"\n[Source {i}: {meta.get('file', 'Unknown')} - Page {meta.get('page', '?')}]\n{c['text']}\n"
        
        lang_note = "Answer in roman (roman )" if language == "hi" else "Answer in clear English"
        
        prompt = f"""You are an AI Document Assistant by {CREATOR}.

{lang_note}

INSTRUCTIONS:
- Answer ONLY from provided context
- Cite sources as [📄 Filename - Page X]
- Be clear and structured
- If not found, say so

CONTEXT:
{context_text}

QUESTION: {question}

ANSWER:"""
        
        answer = self.generate(prompt)
        
        sources = []
        seen = set()
        for c in context[:5]:
            meta = c.get("meta", {})
            key = f"{meta.get('file', '')}-{meta.get('page', '')}"
            if key not in seen and meta.get('file'):
                sources.append({
                    "file": meta.get("file"),
                    "page": meta.get("page")
                })
                seen.add(key)
        
        return {"answer": answer, "sources": sources}
    
    def summarize(self, content, style, language):
        """Summarize content"""
        lang = "in roman" if language == "hi" else "in English"
        
        styles = {
            "brief": f"Write a brief summary (150-200 words) {lang}",
            "detailed": f"Write a comprehensive summary {lang}",
            "bullets": f"Write a bullet-point summary {lang}",
            "cheatsheet": f"Create a quick-reference cheat sheet {lang}"
        }
        
        prompt = f"""{styles.get(style, styles['brief'])}

Content:
{content[:15000]}

Summary:"""
        return self.generate(prompt)
    
    def create_notes(self, content, style, language):
        """Create study notes"""
        lang = "in roman" if language == "hi" else "in English"
        
        styles = {
            "detailed": f"Create comprehensive study notes {lang} with headings, explanations, and examples",
            "revision": f"Create concise revision notes {lang} for quick review",
            "cheatsheet": f"Create an exam cheat sheet {lang} with key facts and formulas"
        }
        
        prompt = f"""{styles.get(style, styles['detailed'])}

Make it student-friendly and easy to understand.

Content:
{content[:15000]}

Notes:"""
        return self.generate(prompt)
    
    def create_mcqs(self, content, count, language):
        """Generate MCQs"""
        lang = "in roman" if language == "hi" else "in English"
        
        prompt = f"""Create {count} high-quality MCQs {lang}.

Format each question as:

Q[N]. [Question]
A) [Option A]
B) [Option B]
C) [Option C]
D) [Option D]

✅ Answer: [Letter]
💡 Explanation: [Brief explanation]

---

Content:
{content[:15000]}

Generate {count} MCQs:"""
        return self.generate(prompt)
    
    def create_flashcards(self, content, count, language):
        """Generate flashcards"""
        lang = "in roman" if language == "hi" else "in English"
        
        prompt = f"""Create {count} study flashcards {lang}.

Format:

🎴 Card [N]
📌 FRONT: [Question/Term]
📖 BACK: [Answer/Definition]
💡 Tip: [Memory hint]

---

Content:
{content[:15000]}

Create {count} flashcards:"""
        return self.generate(prompt)


# ============================================
# 🎯 SMART ANSWER ENGINE (ChatGPT Style)
# ============================================

def smart_answer_engine(question, language="en"):
    """
    Hybrid AI - ALWAYS answers from both documents AND general knowledge
    Kabhi bhi "not found" nahi bolega
    """
    try:
        # Step 1: Documents mein search karo
        doc_results = []
        doc_context = ""
        has_docs = False
        
        if st.session_state.db and st.session_state.docs:
            doc_results = st.session_state.db.search(question, k=5)
            if doc_results and len(doc_results) > 0:
                has_docs = True
                for i, result in enumerate(doc_results[:3], 1):
                    meta = result.get("meta", {})
                    doc_context += f"\n[Source {i}: {meta.get('file', 'Doc')} - Page {meta.get('page', '?')}]\n{result['text']}\n"
        
        # Step 2: ALWAYS use general knowledge + documents (hybrid approach)
        if language == "hi":
            if has_docs:
                # Documents MILے + AI knowledge
                prompt = f"""You are a helpful AI assistant. Answer in ROMAN HINGLISH (Hindi words in English).

**IMPORTANT:** 
- Main ne kuch documents upload kiye hain (niche diye hue hain)
- Tum documents KA info use karo PLUS apni general knowledge bhi add karo
- Agar documents mein poori info nahi hai, to baaki general knowledge se complete karo
- Kabhi "not found" mat kehna, hamesha helpful answer do

**UPLOADED DOCUMENTS:**
{doc_context}

**STUDENT KA QUESTION:** {question}

**TUMHARA COMPLETE ANSWER (Roman Hinglish mein):**

Format (agar documents relevant hain):
📚 **Documents mein kya hai:**
[Jo documents mein mila]

🧠 **Aur bhi info (General Knowledge):**
[Extra helpful information jo documents mein nahi tha]

---
Format (agar documents mein nahi mila):
🧠 **General Knowledge se Answer:**
[Complete answer with examples and explanations]

Note: Documents mein specific info nahi mili, lekin main general knowledge se explain kar raha hoon.

Ab answer do (friendly aur detailed):"""
            else:
                # NO documents - pure ChatGPT mode
                prompt = f"""You are a helpful AI study assistant. Answer in ROMAN HINGLISH.

Koi documents upload nahi hain, lekin koi baat nahi!

**QUESTION:** {question}

Ek detailed aur helpful answer do Roman Hinglish mein:
- Clear explanation
- Examples agar possible ho
- Step-by-step agar complex topic hai
- Friendly tone rakho

Answer:"""
        
        else:  # English mode
            if has_docs:
                # Documents + AI knowledge
                prompt = f"""You are a helpful AI assistant providing comprehensive answers.

**IMPORTANT INSTRUCTIONS:**
- I have some uploaded documents (shown below)
- Use information from documents PLUS your general knowledge
- If documents don't have complete info, fill in with general knowledge
- NEVER say "not found" - always provide a helpful answer

**UPLOADED DOCUMENTS:**
{doc_context}

**STUDENT'S QUESTION:** {question}

**YOUR COMPREHENSIVE ANSWER:**

Format (if documents are relevant):
📚 **From Uploaded Documents:**
[Information found in the documents]

🧠 **Additional Information (General Knowledge):**
[Extra helpful details to complete the answer]

---
Format (if documents don't have the info):
🧠 **Answer from General Knowledge:**
[Complete answer with examples and explanations]

Note: Specific information wasn't in the documents, but I'm providing a comprehensive answer from general knowledge.

Provide your detailed answer:"""
            else:
                # NO documents - pure ChatGPT mode
                prompt = f"""You are a helpful AI study assistant.

No documents are uploaded, but that's okay!

**QUESTION:** {question}

Provide a detailed, helpful answer with:
- Clear explanations
- Examples where relevant
- Step-by-step for complex topics
- Friendly, educational tone

Answer:"""
        
        # Step 3: Generate answer
        answer = st.session_state.ai.generate(prompt)
        
        # Step 4: Extract sources (agar documents se info mili)
        sources = []
        if has_docs:
            seen = set()
            for result in doc_results[:3]:
                meta = result.get("meta", {})
                file = meta.get('file', '')
                page = meta.get('page', '')
                if file and f"{file}-{page}" not in seen:
                    sources.append({"file": file, "page": page})
                    seen.add(f"{file}-{page}")
        
        return {
            "answer": answer,
            "sources": sources,
            "has_docs": has_docs,
            "doc_count": len(doc_results) if doc_results else 0
        }
        
    except Exception as e:
        # Error handling
        error_msg = f"Error occurred: {str(e)[:200]}"
        if language == "hi":
            error_msg = f"Error aaya: {str(e)[:200]}\n\nKoi baat nahi, question phir se poocho!"
        
        return {
            "answer": error_msg,
            "sources": [],
            "has_docs": False,
            "doc_count": 0
        }
# ============================================
# 💾 SESSION STATE
# ============================================

if 'initialized' not in st.session_state:
    st.session_state.initialized = True
    st.session_state.logged_in = False
    st.session_state.username = ""
    st.session_state.docs = {}
    st.session_state.chat = []
    st.session_state.student_chat = []
    st.session_state.db = None
    st.session_state.ai = None
    st.session_state.connected = False
    st.session_state.model = None
    st.session_state.mode = "enterprise"
    st.session_state.lang = "en"
    st.session_state.processor = DocumentProcessor()
    st.session_state.api_key = None
    st.session_state.connection_error = None


# ============================================
# 🔐 LOGIN PAGE
# ============================================

def show_login_page():
    """Display login page"""
    st.markdown('<h1 class="hero-header">✨ AI Document Studio</h1>', unsafe_allow_html=True)
    st.markdown('<p class="hero-sub">Intelligent Document Analysis & Study Assistant</p>', unsafe_allow_html=True)
    st.markdown(f'<p class="hero-creator">Created by {CREATOR}</p>', unsafe_allow_html=True)
    
    st.markdown('<div class="fancy-divider"></div>', unsafe_allow_html=True)
    
    # Login Card
    col1, col2, col3 = st.columns([1, 2, 1])
    
    with col2:
        st.markdown("""
        <div class="login-container">
            <div class="login-title">🔐 Welcome Back</div>
            <div class="login-subtitle">Sign in to continue to your workspace</div>
        </div>
        """, unsafe_allow_html=True)
        
        with st.form("login_form"):
            username = st.text_input("👤 Username", placeholder="Enter your username")
            password = st.text_input("🔒 Password", type="password", placeholder="Enter your password")
            
            col_a, col_b = st.columns(2)
            with col_a:
                remember = st.checkbox("Remember me")
            
            submit = st.form_submit_button("🚀 Sign In", use_container_width=True)
            
            if submit:
                if username and password:
                    # Simple authentication (in production, use proper auth)
                    if len(username) >= 3 and len(password) >= 4:
                        st.session_state.logged_in = True
                        st.session_state.username = username
                        st.rerun()
                    else:
                        show_error("Invalid credentials. Username min 3 chars, Password min 4 chars.")
                else:
                    show_warning("Please enter both username and password.")
        
        st.markdown("<br>", unsafe_allow_html=True)
        
        # Guest Access
        if st.button("👁️ Continue as Guest", use_container_width=True):
            st.session_state.logged_in = True
            st.session_state.username = "Guest"
            st.rerun()
    
    # Features Section
    st.markdown('<div class="fancy-divider"></div>', unsafe_allow_html=True)
    st.markdown('<h2 style="text-align:center; color: white; margin-bottom: 30px;">✨ Features</h2>', unsafe_allow_html=True)
    
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        st.markdown("""
        <div class="feature-box">
            <span class="feature-icon">📄</span>
            <div class="feature-title">Multi-Format</div>
            <div class="feature-desc">PDF, Word, Excel, PowerPoint, Text files</div>
        </div>
        """, unsafe_allow_html=True)
    
    with col2:
        st.markdown("""
        <div class="feature-box">
            <span class="feature-icon">🧠</span>
            <div class="feature-title">Smart AI</div>
            <div class="feature-desc">Advanced semantic search & analysis</div>
        </div>
        """, unsafe_allow_html=True)
    
    with col3:
        st.markdown("""
        <div class="feature-box">
            <span class="feature-icon">📚</span>
            <div class="feature-title">Study Tools</div>
            <div class="feature-desc">Notes, MCQs, Flashcards generator</div>
        </div>
        """, unsafe_allow_html=True)
    
    with col4:
        st.markdown("""
        <div class="feature-box">
            <span class="feature-icon">🌍</span>
            <div class="feature-title">Bilingual</div>
            <div class="feature-desc">English & roman support</div>
        </div>
        """, unsafe_allow_html=True)


# ============================================
# 📱 SIDEBAR
# ============================================

def show_sidebar():
    """Display sidebar"""
    with st.sidebar:
        # Header
        st.markdown(f"""
        <div style="text-align: center; padding: 25px 0; border-bottom: 1px solid rgba(255,255,255,0.1); margin-bottom: 25px;">
            <div style="font-size: 3rem; margin-bottom: 12px;">✨</div>
            <h2 style="color: white; font-size: 1.4rem; font-weight: 700; margin: 0; font-family: 'Space Grotesk', sans-serif;">
                Doc Studio
            </h2>
            <p style="color: rgba(255,255,255,0.5); font-size: 0.85rem; margin: 8px 0 0 0;">
                {VERSION}
            </p>
            <div style="margin-top: 15px;">
                <span class="badge badge-success">👤 {st.session_state.username}</span>
            </div>
        </div>
        """, unsafe_allow_html=True)
        
        # API Connection
        st.markdown("### 🔌 API Connection")
        
        api_key = API_KEY or st.text_input(
            "Gemini API Key",
            type="password",
            placeholder="Enter your API key...",
            help="Get free at: aistudio.google.com"
        )
        
        # Connect button for manual retry
        connect_btn = st.button("🔌 Connect", use_container_width=True, disabled=st.session_state.connected or not api_key)
        
        if api_key and not st.session_state.connected and connect_btn:
            with st.spinner("🔄 Connecting to Gemini AI..."):
                model, error = find_model(api_key)
                if model and error is None:
                    # Configure API globally
                    genai.configure(api_key=api_key.strip())
                    st.session_state.api_key = api_key.strip()
                    st.session_state.model = model
                    st.session_state.db = VectorStore()
                    st.session_state.ai = AIAssistant(model)
                    st.session_state.connected = True
                    st.session_state.connection_error = None
                    st.rerun()
                else:
                    st.session_state.connection_error = error
                    st.rerun()
        
        # Auto-connect on first load if API key exists
        if api_key and not st.session_state.connected and st.session_state.connection_error is None:
            with st.spinner("🔄 Connecting to Gemini AI..."):
                model, error = find_model(api_key)
                if model and error is None:
                    genai.configure(api_key=api_key.strip())
                    st.session_state.api_key = api_key.strip()
                    st.session_state.model = model
                    st.session_state.db = VectorStore()
                    st.session_state.ai = AIAssistant(model)
                    st.session_state.connected = True
                    st.session_state.connection_error = None
                    st.rerun()
                else:
                    st.session_state.connection_error = error
        
        if st.session_state.connected:
            st.markdown(f"""
            <div class="success-box" style="margin: 10px 0;">
                ✅ Connected<br>
                <small style="opacity: 0.8;">Model: {st.session_state.model}</small>
            </div>
            """, unsafe_allow_html=True)
        elif st.session_state.connection_error:
            st.markdown(f'<div class="error-box">❌ {st.session_state.connection_error}</div>', unsafe_allow_html=True)
            if st.button("🔄 Retry Connection", use_container_width=True):
                st.session_state.connection_error = None
                st.rerun()
        
        st.markdown('<div class="fancy-divider"></div>', unsafe_allow_html=True)
        
        # Mode Selection
        st.markdown("### 🎯 Mode")
        mode = st.radio(
            "Select mode",
            ["🏢 Enterprise", "📚 Student"],
            horizontal=True,
            label_visibility="collapsed"
        )
        st.session_state.mode = "enterprise" if "Enterprise" in mode else "student"
        
        st.markdown('<div class="fancy-divider"></div>', unsafe_allow_html=True)
        
        # Document Upload
        if st.session_state.connected:
            st.markdown("### 📁 Documents")
            
            files = st.file_uploader(
                "Upload files",
                type=["pdf", "docx", "xlsx", "pptx", "txt"],
                accept_multiple_files=True,
                label_visibility="collapsed"
            )
            
            if files:
                for file in files:
                    if file.name not in st.session_state.docs:
                        with st.spinner(f"Processing {file.name}..."):
                            try:
                                # 1. Temporary file banana (Vercel RAM ke liye)
                                with tempfile.NamedTemporaryFile(delete=False, suffix=Path(file.name).suffix) as tmp:
                                    tmp.write(file.getvalue())
                                    tmp_path = tmp.name
                                
                                # 2. Document process karna
                                doc = st.session_state.processor.process(tmp_path)
                                
                                # 3. Vector Database mein add karna
                                st.session_state.db.add_document(doc["id"], doc["chunks"])
                                
                                # 4. Session mein save karna
                                st.session_state.docs[file.name] = doc
                                
                                # 5. Temporary file delete karna (Disk clean rakhne ke liye)
                                if os.path.exists(tmp_path):
                                    os.remove(tmp_path)
                                
                                st.success(f"✅ {file.name} successfully processed!")
                                time.sleep(0.3)
                                st.rerun()
                            except Exception as e:
                                st.error(f"❌ Error processing {file.name}: {str(e)}")
            
            # Document List
            if st.session_state.docs:
                st.markdown("#### 📚 Loaded")
                
                for name, doc in list(st.session_state.docs.items()):
                    with st.expander(f"📄 {name[:20]}...", expanded=False):
                        st.caption(f"Type: {doc['type']} | Pages: {doc['pages']}")
                        st.caption(f"Words: {doc['word_count']:,}")
                        
                        if st.button("🗑️ Remove", key=f"rm_{name}", use_container_width=True):
                            st.session_state.db.delete_document(doc['id'])
                            del st.session_state.docs[name]
                            st.rerun()
                
                if st.button("🗑️ Clear All", use_container_width=True, type="secondary"):
                    for doc in st.session_state.docs.values():
                        st.session_state.db.delete_document(doc['id'])
                    st.session_state.docs = {}
                    st.session_state.chat = []
                    st.session_state.student_chat = []
                    st.rerun()
        
        st.markdown('<div class="fancy-divider"></div>', unsafe_allow_html=True)
        
        # Language
        st.markdown("### 🌐 Language")
        lang = st.radio(
            "Response language",
            ["🇬🇧 English", "roman"],
            horizontal=True,
            label_visibility="collapsed"
        )
        st.session_state.lang = "hi" if "roman" in lang else "en"
        
        st.markdown('<div class="fancy-divider"></div>', unsafe_allow_html=True)
        
        # Stats
        st.markdown("### 📊 Stats")
        col1, col2 = st.columns(2)
        with col1:
            st.metric("Docs", len(st.session_state.docs))
        with col2:
            st.metric("Chunks", st.session_state.db.count() if st.session_state.db else 0)
        
        st.markdown('<div class="fancy-divider"></div>', unsafe_allow_html=True)
        
        # Logout
        if st.button("🚪 Logout", use_container_width=True):
            st.session_state.logged_in = False
            st.session_state.username = ""
            st.rerun()


# ============================================
# 🏢 ENTERPRISE MODE
# ============================================

def enterprise_mode():
    """Enterprise mode interface"""
    st.markdown('<h1 class="hero-header">🏢 Enterprise Mode</h1>', unsafe_allow_html=True)
    st.markdown('<p class="hero-sub">Professional Document Intelligence</p>', unsafe_allow_html=True)
    
    tabs = st.tabs(["💬 Chat", "📝 Summarize", "📊 Analytics", "🔧 Tools"])
    
    # CHAT TAB
    with tabs[0]:
        if not st.session_state.docs:
            show_info("📁 Upload documents using the sidebar to start chatting with your files!")
        else:
            st.markdown('<div class="glass-card">', unsafe_allow_html=True)
            
            # Chat History
            for msg in st.session_state.chat:
                if msg["role"] == "user":
                    st.markdown(f'<div class="user-message"><strong>You:</strong> {msg["content"]}</div>', unsafe_allow_html=True)
                else:
                    st.markdown(f'<div class="ai-message"><strong>AI Assistant:</strong><br><br>{msg["content"]}</div>', unsafe_allow_html=True)
                    
                    if msg.get("sources"):
                        sources_html = ""
                        for src in msg["sources"][:3]:
                            sources_html += f'<span class="source-tag">📄 {src["file"]} - Page {src["page"]}</span>'
                        st.markdown(sources_html, unsafe_allow_html=True)
            
            st.markdown('<div class="fancy-divider"></div>', unsafe_allow_html=True)
            
            # Input
            col1, col2 = st.columns([5, 1])
            
            with col1:
                question = st.text_input(
                    "Question",
                    placeholder="Ask anything about your documents...",
                    label_visibility="collapsed",
                    key="ent_q"
                )
            
            with col2:
                ask = st.button("🚀 Ask", use_container_width=True, key="ent_ask")
            
            if ask and question:
                st.session_state.chat.append({"role": "user", "content": question})
                
                with st.spinner("🔍 Analyzing documents..."):
                    results = st.session_state.db.search(question, k=5)
                    
                    if not results:
                        response = {"answer": "No relevant information found in documents.", "sources": []}
                    else:
                        response = st.session_state.ai.answer_question(
                            question, results, st.session_state.lang
                        )
                    
                    st.session_state.chat.append({
                        "role": "assistant",
                        "content": response["answer"],
                        "sources": response.get("sources", [])
                    })
                
                st.rerun()
            
            if st.session_state.chat:
                if st.button("🗑️ Clear Chat", key="ent_clear"):
                    st.session_state.chat = []
                    st.rerun()
            
            st.markdown('</div>', unsafe_allow_html=True)
    
    # SUMMARIZE TAB
    with tabs[1]:
        if not st.session_state.docs:
            show_info("📁 Upload documents first to generate summaries!")
        else:
            st.markdown('<div class="glass-card">', unsafe_allow_html=True)
            
            col1, col2 = st.columns(2)
            
            with col1:
                doc_name = st.selectbox("📄 Select Document", list(st.session_state.docs.keys()), key="sum_doc")
            
            with col2:
                style = st.selectbox("📝 Style", ["Brief", "Detailed", "Bullet Points", "Cheat Sheet"], key="sum_style")
            
            if st.button("✨ Generate Summary", use_container_width=True, key="sum_btn"):
                with st.spinner("Creating summary..."):
                    style_map = {"Brief": "brief", "Detailed": "detailed", "Bullet Points": "bullets", "Cheat Sheet": "cheatsheet"}
                    content = st.session_state.docs[doc_name]["content"]
                    summary = st.session_state.ai.summarize(content, style_map[style], st.session_state.lang)
                    
                    st.markdown("### 📋 Summary")
                    st.markdown(f'<div class="content-box">{summary}</div>', unsafe_allow_html=True)
                    
                    st.download_button(
                        "📥 Download Summary",
                        add_watermark(summary),
                        f"summary_{doc_name}.txt",
                        mime="text/plain",
                        key="sum_dl"
                    )
            
            st.markdown('</div>', unsafe_allow_html=True)
    
    # ANALYTICS TAB
    with tabs[2]:
        if not st.session_state.docs:
            show_info("📁 Upload documents to view analytics!")
        else:
            # Metrics
            col1, col2, col3, col4 = st.columns(4)
            
            total_pages = sum(d["pages"] for d in st.session_state.docs.values())
            total_chunks = sum(len(d["chunks"]) for d in st.session_state.docs.values())
            total_words = sum(d["word_count"] for d in st.session_state.docs.values())
            
            with col1:
                st.markdown(f'<div class="metric-box"><div class="metric-number">{len(st.session_state.docs)}</div><div class="metric-title">Documents</div></div>', unsafe_allow_html=True)
            with col2:
                st.markdown(f'<div class="metric-box"><div class="metric-number">{total_pages}</div><div class="metric-title">Pages</div></div>', unsafe_allow_html=True)
            with col3:
                st.markdown(f'<div class="metric-box"><div class="metric-number">{total_chunks}</div><div class="metric-title">Chunks</div></div>', unsafe_allow_html=True)
            with col4:
                st.markdown(f'<div class="metric-box"><div class="metric-number">{total_words:,}</div><div class="metric-title">Words</div></div>', unsafe_allow_html=True)
            
            st.markdown("<br>", unsafe_allow_html=True)
            
            # Charts
            col1, col2 = st.columns(2)
            
            with col1:
                types = [d["type"] for d in st.session_state.docs.values()]
                type_counts = pd.Series(types).value_counts()
                
                fig = go.Figure(data=[go.Pie(
                    labels=type_counts.index,
                    values=type_counts.values,
                    hole=0.5,
                    marker=dict(colors=['#667eea', '#f093fb', '#4facfe', '#43e97b'])
                )])
                fig.update_layout(
                    title="Document Types",
                    paper_bgcolor='rgba(0,0,0,0)',
                    plot_bgcolor='rgba(0,0,0,0)',
                    font=dict(color='white'),
                    height=300
                )
                st.plotly_chart(fig, use_container_width=True)
            
            with col2:
                names = [n[:12] for n in st.session_state.docs.keys()]
                chunks = [len(d["chunks"]) for d in st.session_state.docs.values()]
                
                fig = go.Figure(data=[go.Bar(
                    x=names,
                    y=chunks,
                    marker=dict(color=chunks, colorscale='Viridis')
                )])
                fig.update_layout(
                    title="Chunks per Document",
                    paper_bgcolor='rgba(0,0,0,0)',
                    plot_bgcolor='rgba(0,0,0,0)',
                    font=dict(color='white'),
                    height=300
                )
                st.plotly_chart(fig, use_container_width=True)
    
    # TOOLS TAB
    with tabs[3]:
        if not st.session_state.docs:
            show_info("📁 Upload documents first!")
        else:
            st.markdown('<div class="glass-card">', unsafe_allow_html=True)
            
            tool = st.selectbox("🔧 Select Tool", ["Keyword Extraction", "Content Analysis", "Document Comparison"])
            
            if tool == "Keyword Extraction":
                doc = st.selectbox("Document", list(st.session_state.docs.keys()), key="kw_doc")
                num = st.slider("Number of keywords", 5, 30, 15)
                
                if st.button("🔍 Extract Keywords", key="kw_btn"):
                    with st.spinner("Extracting..."):
                        prompt = f"Extract top {num} keywords and key concepts from:\n\n{st.session_state.docs[doc]['content'][:10000]}\n\nFormat as a list with brief explanations."
                        result = st.session_state.ai.generate(prompt)
                        st.markdown(f'<div class="content-box">{result}</div>', unsafe_allow_html=True)
            
            elif tool == "Content Analysis":
                doc = st.selectbox("Document", list(st.session_state.docs.keys()), key="an_doc")
                analysis = st.selectbox("Analysis Type", ["Sentiment", "Topics", "Readability"])
                
                if st.button("📊 Analyze", key="an_btn"):
                    with st.spinner("Analyzing..."):
                        prompts = {
                            "Sentiment": "Analyze the sentiment, tone, and emotional content:",
                            "Topics": "Identify main topics and themes:",
                            "Readability": "Assess readability level and target audience:"
                        }
                        prompt = f"{prompts[analysis]}\n\n{st.session_state.docs[doc]['content'][:10000]}"
                        result = st.session_state.ai.generate(prompt)
                        st.markdown(f'<div class="content-box">{result}</div>', unsafe_allow_html=True)
            
            elif tool == "Document Comparison":
                if len(st.session_state.docs) < 2:
                    show_warning("Upload at least 2 documents to compare!")
                else:
                    doc1 = st.selectbox("First Document", list(st.session_state.docs.keys()), key="cmp1")
                    doc2 = st.selectbox("Second Document", list(st.session_state.docs.keys()), key="cmp2")
                    
                    if st.button("🔍 Compare", key="cmp_btn"):
                        with st.spinner("Comparing..."):
                            prompt = f"""Compare these documents:

Document 1: {doc1}
{st.session_state.docs[doc1]['content'][:5000]}

Document 2: {doc2}
{st.session_state.docs[doc2]['content'][:5000]}

Provide: similarities, differences, unique insights from each."""
                            result = st.session_state.ai.generate(prompt)
                            st.markdown(f'<div class="content-box">{result}</div>', unsafe_allow_html=True)
            
            st.markdown('</div>', unsafe_allow_html=True)


# ============================================
# 📚 STUDENT MODE
# ============================================

def student_mode():
    """Student mode interface"""
    st.markdown('<h1 class="hero-header">📚 Student Mode</h1>', unsafe_allow_html=True)
    st.markdown('<p class="hero-sub">Smart Study Assistant & Exam Prep</p>', unsafe_allow_html=True)
    
    tabs = st.tabs(["💬 bossGPT", "📝 Notes", "❓ MCQs", "🎴 Flashcards", "🎯 Exam Prep"])
    
    # ASK TAB
    # ASK TAB - bossGPT (ChatGPT Style)
    with tabs[0]:
        st.markdown('<div class="glass-card">', unsafe_allow_html=True)
        
        # Smart Info Banner
        if st.session_state.docs:
            count = len(st.session_state.docs)
            if st.session_state.lang == "hi":
                st.markdown(f'<div class="info-box">💡 {count} documents loaded! Main documents + AI knowledge DONO use karunga!</div>', unsafe_allow_html=True)
            else:
                st.markdown(f'<div class="info-box">💡 {count} documents loaded! Using BOTH documents & AI knowledge!</div>', unsafe_allow_html=True)
        else:
            if st.session_state.lang == "hi":
                st.markdown('<div class="info-box">🤖 Documents nahi hain? Koi problem nahi! Pure ChatGPT mode mein kaam karunga!</div>', unsafe_allow_html=True)
            else:
                st.markdown('<div class="info-box">🤖 documents ni ha no problem boss hena!</div>', unsafe_allow_html=True)
        
        st.markdown('<div class="fancy-divider"></div>', unsafe_allow_html=True)
        
        # Chat History with Smart Badges
        for msg in st.session_state.student_chat:
            if msg["role"] == "user":
                st.markdown(f'<div class="user-message"><strong>You:</strong> {msg["content"]}</div>', unsafe_allow_html=True)
            else:
                # Smart badge showing source
                if msg.get("has_docs"):
                    badge_text = f'📚 Docs + 🧠 AI ({msg.get("doc_count", 0)} sources)'
                else:
                    badge_text = '🧠 AI Knowledge (BossGPT Mode)'
                
                st.markdown(f'<span class="badge badge-success">{badge_text}</span>', unsafe_allow_html=True)
                st.markdown(f'<div class="ai-message"><strong>bossGPT:</strong><br><br>{msg["content"]}</div>', unsafe_allow_html=True)
                
                # Show sources if available
                if msg.get("sources"):
                    sources_html = '<div style="margin-top: 15px;">'
                    sources_html += '<strong>📖 Sources:</strong><br>'
                    for src in msg["sources"]:
                        sources_html += f'<span class="source-tag">📄 {src["file"]} - Page {src["page"]}</span>'
                    sources_html += '</div>'
                    st.markdown(sources_html, unsafe_allow_html=True)
        
        st.markdown('<div class="fancy-divider"></div>', unsafe_allow_html=True)
        
        # Quick Suggestions
        if not st.session_state.student_chat:
            if st.session_state.lang == "hi":
                st.markdown("#### 💡 Examples:")
                suggestions = ["Main concepts explain karo", "Key formulas kya hain?", "Quick summary do"]
            else:
                st.markdown("#### 💡 Try asking:")
                suggestions = ["Explain main concepts", "What are key formulas?", "Give quick summary"]
            
            cols = st.columns(3)
            for i, col in enumerate(cols):
                with col:
                    if st.button(suggestions[i], key=f"sug_{i}", use_container_width=True):
                        st.session_state.student_chat.append({"role": "user", "content": suggestions[i]})
                        
                        with st.spinner("Thinking..." if st.session_state.lang == "en" else "Soch raha hoon..."):
                            response = smart_answer_engine(suggestions[i], st.session_state.lang)
                            
                            st.session_state.student_chat.append({
                                "role": "assistant",
                                "content": response["answer"],
                                "sources": response.get("sources", []),
                                "has_docs": response.get("has_docs", False),
                                "doc_count": response.get("doc_count", 0)
                            })
                        st.rerun()
        
        # Question Input
        if st.session_state.lang == "hi":
            placeholder = "Kuch bhi poocho... (Documents ya general knowledge, dono chalega!)"
        else:
            placeholder = "Ask me anything... (Documents or general knowledge!)"
        
        question = st.text_input("Your question", placeholder=placeholder, label_visibility="collapsed", key="stu_q")
        
        # Buttons Row
        col1, col2 = st.columns([4, 1])
        
        with col2:
            if st.session_state.student_chat:
                if st.button("🗑️", key="stu_clear", use_container_width=True):
                    st.session_state.student_chat = []
                    st.rerun()
        
        # Ask Button with Smart Engine
        ask_btn_text = "🚀 Poocho" if st.session_state.lang == "hi" else "🚀 Ask"
        
        if st.button(ask_btn_text, key="stu_ask", use_container_width=True):
            if question and question.strip():
                # Add user message
                st.session_state.student_chat.append({"role": "user", "content": question})
                
                # Get smart answer using the engine
                with st.spinner("Analyzing..." if st.session_state.lang == "en" else "Analyze kar raha hoon..."):
                    response = smart_answer_engine(question, st.session_state.lang)
                    
                    st.session_state.student_chat.append({
                        "role": "assistant",
                        "content": response["answer"],
                        "sources": response.get("sources", []),
                        "has_docs": response.get("has_docs", False),
                        "doc_count": response.get("doc_count", 0)
                    })
                
                st.rerun()
            else:
                if st.session_state.lang == "hi":
                    show_warning("Pehle kuch type karo!")
                else:
                    show_warning("Please type a question first!")
        
        st.markdown('</div>', unsafe_allow_html=True)
    
    # NOTES TAB
    with tabs[1]:
        if not st.session_state.docs:
            show_info("📁 Upload your study materials first!")
        else:
            st.markdown('<div class="glass-card">', unsafe_allow_html=True)
            
            col1, col2 = st.columns(2)
            
            with col1:
                doc = st.selectbox("📄 Select Material", list(st.session_state.docs.keys()), key="notes_doc")
            
            with col2:
                style = st.selectbox("📝 Notes Style", ["Detailed Notes", "Quick Revision", "Cheat Sheet"], key="notes_style")
            
            if st.button("✨ Generate Notes", use_container_width=True, key="notes_btn"):
                with st.spinner("Creating your notes..."):
                    style_map = {"Detailed Notes": "detailed", "Quick Revision": "revision", "Cheat Sheet": "cheatsheet"}
                    content = st.session_state.docs[doc]["content"]
                    notes = st.session_state.ai.create_notes(content, style_map[style], st.session_state.lang)
                    
                    st.markdown("### 📚 Your Study Notes")
                    st.markdown(f'<div class="content-box">{notes}</div>', unsafe_allow_html=True)
                    
                    st.download_button(
                        "📥 Download Notes",
                        add_watermark(notes),
                        f"notes_{doc}.txt",
                        mime="text/plain",
                        key="notes_dl"
                    )
            
            st.markdown('</div>', unsafe_allow_html=True)
    
    # MCQ TAB
    with tabs[2]:
        if not st.session_state.docs:
            show_info("📁 Upload your study materials first!")
        else:
            st.markdown('<div class="glass-card">', unsafe_allow_html=True)
            
            col1, col2 = st.columns(2)
            
            with col1:
                doc = st.selectbox("📄 Select Material", list(st.session_state.docs.keys()), key="mcq_doc")
            
            with col2:
                num = st.slider("Number of Questions", 5, 25, 10, key="mcq_num")
            
            if st.button("🎯 Generate Quiz", use_container_width=True, key="mcq_btn"):
                with st.spinner("Creating your quiz..."):
                    content = st.session_state.docs[doc]["content"]
                    mcqs = st.session_state.ai.create_mcqs(content, num, st.session_state.lang)
                    
                    st.markdown("### ❓ Practice Quiz")
                    st.markdown(f'<div class="content-box">{mcqs}</div>', unsafe_allow_html=True)
                    
                    st.download_button(
                        "📥 Download Quiz",
                        add_watermark(mcqs),
                        f"quiz_{doc}.txt",
                        mime="text/plain",
                        key="mcq_dl"
                    )
            
            st.markdown('</div>', unsafe_allow_html=True)
    
    # FLASHCARDS TAB
    with tabs[3]:
        if not st.session_state.docs:
            show_info("📁 Upload your study materials first!")
        else:
            st.markdown('<div class="glass-card">', unsafe_allow_html=True)
            
            col1, col2 = st.columns(2)
            
            with col1:
                doc = st.selectbox("📄 Select Material", list(st.session_state.docs.keys()), key="fc_doc")
            
            with col2:
                num = st.slider("Number of Cards", 10, 40, 20, key="fc_num")
            
            if st.button("🎴 Generate Flashcards", use_container_width=True, key="fc_btn"):
                with st.spinner("Creating flashcards..."):
                    content = st.session_state.docs[doc]["content"]
                    cards = st.session_state.ai.create_flashcards(content, num, st.session_state.lang)
                    
                    st.markdown("### 🎴 Your Flashcards")
                    st.markdown(f'<div class="content-box">{cards}</div>', unsafe_allow_html=True)
                    
                    st.download_button(
                        "📥 Download Flashcards",
                        add_watermark(cards),
                        f"flashcards_{doc}.txt",
                        mime="text/plain",
                        key="fc_dl"
                    )
            
            st.markdown('</div>', unsafe_allow_html=True)
    
    # EXAM PREP TAB
    with tabs[4]:
        if not st.session_state.docs:
            show_info("📁 Upload your study materials first!")
        else:
            st.markdown('<div class="glass-card">', unsafe_allow_html=True)
            
            doc = st.selectbox("📄 Select Material", list(st.session_state.docs.keys()), key="prep_doc")
            
            tool = st.selectbox("🎯 Prep Tool", [
                "📋 Study Plan",
                "🎯 Important Topics",
                "⏰ Quick Revision",
                "🔮 Predict Questions"
            ], key="prep_tool")
            
            if "Study Plan" in tool:
                days = st.number_input("Days until exam", 1, 60, 7, key="prep_days")
                
                if st.button("📅 Create Plan", key="prep_btn"):
                    with st.spinner("Creating your plan..."):
                        prompt = f"""Create a {days}-day study plan for this material:

{st.session_state.docs[doc]['content'][:10000]}

Include daily schedule, topics, and revision time."""
                        result = st.session_state.ai.generate(prompt)
                        st.markdown(f'<div class="content-box">{result}</div>', unsafe_allow_html=True)
                        
                        st.download_button("📥 Download Plan", add_watermark(result), "study_plan.txt", key="plan_dl")
            
            elif "Important Topics" in tool:
                if st.button("🎯 Find Topics", key="topic_btn"):
                    with st.spinner("Analyzing..."):
                        prompt = f"""Identify most important exam topics from:

{st.session_state.docs[doc]['content'][:12000]}

Rank by importance with explanations."""
                        result = st.session_state.ai.generate(prompt)
                        st.markdown(f'<div class="content-box">{result}</div>', unsafe_allow_html=True)
            
            elif "Quick Revision" in tool:
                if st.button("⚡ Quick Revision", key="rev_btn"):
                    with st.spinner("Creating..."):
                        prompt = f"""Create a 15-minute last-minute revision guide for:

{st.session_state.docs[doc]['content'][:10000]}

Include key points, formulas, and must-remember facts."""
                        result = st.session_state.ai.generate(prompt)
                        st.markdown(f'<div class="content-box">{result}</div>', unsafe_allow_html=True)
                        
                        st.download_button("📥 Download", add_watermark(result), "quick_revision.txt", key="rev_dl")
            
            elif "Predict Questions" in tool:
                num = st.slider("Questions to predict", 5, 20, 10, key="pred_num")
                
                if st.button("🔮 Predict", key="pred_btn"):
                    with st.spinner("Predicting..."):
                        prompt = f"""Predict {num} most likely exam questions from:

{st.session_state.docs[doc]['content'][:12000]}

Include question, difficulty level, and key answer points."""
                        result = st.session_state.ai.generate(prompt)
                        st.markdown(f'<div class="content-box">{result}</div>', unsafe_allow_html=True)
                        
                        st.download_button("📥 Download", add_watermark(result), "predictions.txt", key="pred_dl")
            
            st.markdown('</div>', unsafe_allow_html=True)


# ============================================
# 📄 FOOTER
# ============================================

def show_footer():
    """Display footer"""
    st.markdown('<div class="fancy-divider"></div>', unsafe_allow_html=True)
    
    st.markdown(f"""
    <div class="footer-container">
        <div class="footer-title">✨ AI Document Studio</div>
        <p style="color: rgba(255,255,255,0.6);">
            Model: <strong style="color: white;">{st.session_state.model or 'Not Connected'}</strong> | 
            Version: <strong style="color: white;">{VERSION}</strong>
        </p>
        <div class="footer-creator">Created by {CREATOR}</div>
        <p class="footer-text">
            Powered by Google Gemini AI & Streamlit<br>
            <span style="font-size: 0.85rem;">© {datetime.now().year} All Rights Reserved</span>
        </p>
    </div>
    """, unsafe_allow_html=True)


# ============================================
# 🚀 MAIN APP
# ============================================

def main():
    """Main application"""
    
    # Check login status
    if not st.session_state.logged_in:
        show_login_page()
        return
    
    # Show sidebar
    show_sidebar()
    
    # Check connection
    if not st.session_state.connected:
        st.markdown('<h1 class="hero-header">✨ AI Document Studio</h1>', unsafe_allow_html=True)
        st.markdown('<p class="hero-sub">Intelligent Document Analysis & Study Assistant</p>', unsafe_allow_html=True)
        st.markdown(f'<p class="hero-creator">Welcome, {st.session_state.username}! 👋</p>', unsafe_allow_html=True)
        
        st.markdown('<div class="fancy-divider"></div>', unsafe_allow_html=True)
        
        show_info("🔌 Please enter your Gemini API key in the sidebar to get started. Get a free key at: aistudio.google.com")
        
        # Features
        st.markdown("### ✨ What you can do")
        
        col1, col2, col3, col4 = st.columns(4)
        
        features = [
            ("📄", "Multi-Format Support", "PDF, Word, Excel, PowerPoint"),
            ("🧠", "AI-Powered Analysis", "Smart Q&A & Insights"),
            ("📚", "Study Tools", "Notes, MCQs, Flashcards"),
            ("🎯", "Exam Prep", "Study plans & predictions")
        ]
        
        for i, (icon, title, desc) in enumerate(features):
            with [col1, col2, col3, col4][i]:
                st.markdown(f"""
                <div class="feature-box">
                    <span class="feature-icon">{icon}</span>
                    <div class="feature-title">{title}</div>
                    <div class="feature-desc">{desc}</div>
                </div>
                """, unsafe_allow_html=True)
        
        show_footer()
        return
    
    # Show mode based on selection
    if st.session_state.mode == "enterprise":
        enterprise_mode()
    else:
        student_mode()
    
    # Footer
    show_footer()


# Run the app
if __name__ == "__main__":
    main()